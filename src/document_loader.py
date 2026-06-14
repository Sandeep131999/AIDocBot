"""
DOCUMENT LOADER & CHUNKING R&D - MARKDOWN-FIRST PIPELINE
========================================================
Any uploaded file format is converted to Markdown for better readability
before chunking and embedding. This preserves structure, headers, tables,
and formatting in a LLM-friendly format.

All settings read from .env file (UTF-8 encoding).

Strategies:
1. recursive: RecursiveCharacterTextSplitter (default, best for most)
2. semantic: SemanticChunker (groups by meaning, best for accuracy)
3. token: TokenTextSplitter (token-aware, best for LLM alignment)
4. sentence: Sentence-based (preserves complete thoughts)
5. markdown: MarkdownHeaderTextSplitter (preserves headers)
"""

import json
import os
import re
import shutil
import subprocess
import nltk
from pathlib import Path
from typing import Any, List, Union, Optional
from langchain_community.document_loaders import (
    PyPDFLoader, TextLoader, CSVLoader,
    UnstructuredExcelLoader, Docx2txtLoader
)
from langchain_text_splitters import (
    RecursiveCharacterTextSplitter,
    TokenTextSplitter,
    MarkdownHeaderTextSplitter
)
from langchain_experimental.text_splitter import SemanticChunker
from langchain_core.documents import Document
from langchain_huggingface import HuggingFaceEmbeddings

from src.config import Config

# Download NLTK data
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt', quiet=True)


# ============================================================================
# STRATEGY NAME NORMALIZATION
# ============================================================================

STRATEGY_ALIASES = {
    "recursive": "recursive",
    "semantic": "semantic",
    "token": "token",
    "sentence": "sentence",
    "markdown": "markdown",
    "markdown-aware": "markdown",
    "markdown_aware": "markdown",
    "markdown aware": "markdown",
    "Markdown-Aware Chunking": "markdown",
    "Markdown-Aware": "markdown",
    "header": "markdown",
    "headers": "markdown",
    "recursive character": "recursive",
    "recursive_character": "recursive",
    "character": "recursive",
    "default": "recursive",
    "semantic chunking": "semantic",
    "semantic_chunking": "semantic",
    "token-based": "token",
    "token_based": "token",
    "tiktoken": "token",
    "sentence-based": "sentence",
    "sentence_based": "sentence",
    "nltk": "sentence",
}


def normalize_strategy(strategy: str) -> str:
    """Normalize any strategy name/alias to canonical form."""
    if not strategy:
        return "recursive"
    normalized = strategy.strip()
    if normalized in STRATEGY_ALIASES:
        return STRATEGY_ALIASES[normalized]
    lower = normalized.lower()
    if lower in STRATEGY_ALIASES:
        return STRATEGY_ALIASES[lower]
    for alias, canonical in STRATEGY_ALIASES.items():
        if lower in alias.lower() or alias.lower() in lower:
            return canonical
    print(f"   ⚠️  Unknown strategy '{strategy}', falling back to 'recursive'")
    return "recursive"


# ============================================================================
# SMART JSON TO MARKDOWN CONVERTER
# ============================================================================

class JSONToMarkdownConverter:
    """Converts JSON knowledge base records to rich, structured Markdown."""

    # Patterns for content type detection
    DATE_PATTERN = re.compile(r'^(\d{1,2}[-/][A-Za-z]{3}[-/]\d{4}|\d{2}-\d{2}-\d{4}|\d{4}-\d{2}-\d{2})$')
    SERIAL_NO_PATTERN = re.compile(r'^\d{1,2}$')
    EMP_ID_PATTERN = re.compile(r'^\d{5,6}$')

    def convert(self, record: dict) -> str:
        """Convert a single JSON record to Markdown."""
        label = record.get('label', '')
        title = record.get('title', '')
        url = record.get('url', '')
        content = record.get('content', '')

        # Extract clean page title
        page_title = title.split('|')[-1].strip() if '|' in title else title
        if not page_title:
            page_title = label

        md_lines = [
            f"# {page_title}",
            "",
            f"**Source:** [{url}]({url})  ",
            f"**Label:** `{label}`",
            "",
            "---",
            "",
        ]

        # Clean noise lines
        content = self._clean_content(content)
        lines = content.split('\n')

        # Detect and format by content type
        content_type = self._detect_content_type(lines, content)

        if content_type == "employee_roster":
            md_lines.extend(self._format_employee_roster(lines, page_title))
        elif content_type == "employee_cards":
            md_lines.extend(self._format_employee_cards(lines))
        elif content_type == "event_list":
            md_lines.extend(self._format_event_list(lines))
        elif content_type == "speech_schedule":
            md_lines.extend(self._format_speech_schedule(lines))
        elif content_type == "tabular_data":
            md_lines.extend(self._format_tabular_data(lines))
        elif content_type == "holiday_list":
            md_lines.extend(self._format_holiday_list(lines))
        elif content_type == "colon_bullets":
            md_lines.extend(self._format_colon_bullets(lines))
        else:
            md_lines.extend(self._format_plain_text(lines))

        return "\n".join(md_lines)

    def convert_batch(self, records: List[dict]) -> str:
        """Convert multiple records to one combined Markdown document."""
        parts = []
        for record in records:
            parts.append(self.convert(record))
            parts.append("\n\n---\n\n")
        return "\n".join(parts)

    def _clean_content(self, content: str) -> str:
        """Remove noise lines from scraped content."""
        content = re.sub(r'^No language is assigned to this page\.?\n?', '', content, flags=re.MULTILINE)
        content = re.sub(r'^English \(en\)\n?', '', content, flags=re.MULTILINE)
        content = re.sub(r'^Home\s*Page\s*$', '', content, flags=re.MULTILINE)
        content = re.sub(r'^Back\s*$', '', content, flags=re.MULTILINE)
        return content.strip()

    def _detect_content_type(self, lines: List[str], content: str) -> str:
        """Auto-detect the structure of the content."""
        stripped = [l.strip() for l in lines if l.strip()]

        # Employee roster: Serial Number + Employee ID + Date Of Joining headers
        if ('Serial Number' in content and 'Employee ID' in content
                and 'Date Of Joining' in content):
            return "employee_roster"

        # Employee cards: repeating Name:/Employee ID:/Designation: pattern
        name_colon_count = sum(1 for l in stripped if l == 'Name:')
        if name_colon_count >= 2 and 'Designation:' in content:
            if 'About Him:' in content or 'About Her:' in content:
                return "employee_cards"
            # Could be freshers list without About sections
            if 'Date Of Joining:' in content:
                return "employee_cards"

        # Event list: Date/Description headers with date entries
        if 'Date' in content and 'Description' in content:
            date_matches = sum(1 for l in stripped if self.DATE_PATTERN.match(l))
            if date_matches >= 2:
                return "event_list"

        # Speech schedule: DATE/EMPLOYEE ID/EMPLOYEE NAME/TOPICS headers
        if ('DATE' in content or 'Date' in content) and 'EMPLOYEE ID' in content and 'TOPICS' in content:
            return "speech_schedule"

        # Holiday list: SL No./Date/Day/Name headers
        if ('SL No.' in content or 'SL No' in content) and 'Day Of The Week' in content:
            return "holiday_list"

        # Tabular data with Name/Date/Visitor/Description
        if ('Name' in stripped[:10] and 'Visited Date' in content
                and 'Visitor' in content and 'Description' in content):
            return "tabular_data"

        # Colon-separated key-value pairs
        colon_lines = [l for l in stripped if ':' in l and not l.startswith('http')]
        if len(colon_lines) >= 3 and len(stripped) > 0 and len(colon_lines) / len(stripped) > 0.3:
            return "colon_bullets"

        return "plain_text"

    # ------------------------------------------------------------------
    # Formatters
    # ------------------------------------------------------------------

    def _format_employee_roster(self, lines: List[str], page_title: str) -> List[str]:
        """Format employee roster as Markdown table."""
        md = [f"## {page_title}", ""]

        headers = ['S.No', 'Employee ID', 'Name', 'Designation', 'Email', 'Date of Joining']
        md.append('| ' + ' | '.join(headers) + ' |')
        md.append('|' + '|'.join(['---'] * len(headers)) + '|')

        i = 0
        while i < len(lines):
            line = lines[i].strip()
            if self.SERIAL_NO_PATTERN.match(line) and len(line) <= 2:
                sno = line
                emp_id = lines[i+1].strip() if i+1 < len(lines) else ''
                doj = lines[i+2].strip() if i+2 < len(lines) else ''

                name = designation = email = ''
                j = i + 3
                while j < len(lines):
                    l = lines[j].strip()
                    if self.SERIAL_NO_PATTERN.match(l) and len(l) <= 2:
                        break
                    if l == 'Name:':
                        name = lines[j+1].strip() if j+1 < len(lines) else ''
                    elif l == 'Designation:':
                        designation = lines[j+1].strip() if j+1 < len(lines) else ''
                    elif l in ('Email ID:', 'Email Id:'):
                        email = lines[j+1].strip() if j+1 < len(lines) else ''
                    j += 1

                md.append(f"| {sno} | {emp_id} | {name} | {designation} | {email} | {doj} |")
                i = j
            else:
                i += 1

        return md

    def _format_employee_cards(self, lines: List[str]) -> List[str]:
        """Format employee cards as structured sections with headers."""
        md = []
        current = {}
        i = 0

        while i < len(lines):
            line = lines[i].strip()

            if line == 'Name:':
                if current:
                    md.extend(self._person_card_to_md(current))
                    md.append("")
                current = {'Name': lines[i+1].strip() if i+1 < len(lines) else ''}
                i += 2
            elif line.rstrip(':') in ('Employee ID', 'Designation', 'Date Of Joining',
                                       'Email Id', 'Email ID', 'Date of Joining'):
                key = line.rstrip(':')
                val = lines[i+1].strip() if i+1 < len(lines) else ''
                current[key] = val
                i += 2
            elif line.startswith('About '):
                key = line.rstrip(':')
                about_lines = []
                j = i + 1
                while j < len(lines):
                    nl = lines[j].strip()
                    if nl == 'Name:':
                        break
                    about_lines.append(nl)
                    j += 1
                current[key] = ' '.join(filter(None, about_lines))
                i = j
            else:
                i += 1

        if current:
            md.extend(self._person_card_to_md(current))

        return md

    def _person_card_to_md(self, person: dict) -> List[str]:
        """Convert person dict to Markdown card."""
        name = person.get('Name', 'Unknown')
        md = [f"### {name}", ""]

        field_order = ['Employee ID', 'Designation', 'Date Of Joining',
                       'Date of Joining', 'Email Id', 'Email ID']
        for key in field_order:
            if key in person and person[key]:
                md.append(f"- **{key}:** {person[key]}")

        # About sections
        for key, val in person.items():
            if key.startswith('About ') and val:
                md.append(f"- **{key}:** {val}")

        return md

    def _format_event_list(self, lines: List[str]) -> List[str]:
        """Format event list with dates as structured sections."""
        md = ["## Events & Activities", ""]
        current = {}
        i = 0

        while i < len(lines):
            line = lines[i].strip()

            if self.DATE_PATTERN.match(line):
                if current:
                    md.extend(self._event_to_md(current))
                    md.append("")
                current = {'Date': line}
                i += 1
            elif line in ('Event:', 'Description:', 'Activity:'):
                key = line.rstrip(':')
                val_lines = []
                j = i + 1
                while j < len(lines):
                    nl = lines[j].strip()
                    if nl in ('Event:', 'Description:', 'Activity:', 'Check the photo gallery here.'):
                        if nl == 'Check the photo gallery here.':
                            val_lines.append('[Photo Gallery]')
                        break
                    if self.DATE_PATTERN.match(nl):
                        break
                    val_lines.append(nl)
                    j += 1
                current[key] = ' '.join(filter(None, val_lines))
                i = j
            else:
                i += 1

        if current:
            md.extend(self._event_to_md(current))

        return md

    def _event_to_md(self, event: dict) -> List[str]:
        """Convert event dict to Markdown."""
        date = event.get('Date', '')
        name = event.get('Event', '')
        md = [f"### {name} ({date})"]
        if 'Description' in event and event['Description']:
            md.append(event['Description'])
        if 'Activity' in event and event['Activity']:
            md.append(f"- **Activity:** {event['Activity']}")
        return md

    def _format_speech_schedule(self, lines: List[str]) -> List[str]:
        """Format morning speech schedule as Markdown table."""
        md = ["## Morning Speech Schedule", ""]
        md.append("| Date | Employee ID | Employee Name | Topic |")
        md.append("|------|-------------|---------------|-------|")

        i = 0
        while i < len(lines):
            line = lines[i].strip()
            if re.match(r'^\d{2}-\d{2}-\d{4}$', line):
                date = line
                emp_id = lines[i+1].strip() if i+1 < len(lines) else ''
                emp_name = lines[i+2].strip() if i+2 < len(lines) else ''
                topic = lines[i+3].strip() if i+3 < len(lines) else ''
                md.append(f"| {date} | {emp_id} | {emp_name} | {topic} |")
                i += 4
            else:
                i += 1

        # Add note if present
        note_lines = [l for l in lines if l.strip().lower().startswith('note:')]
        if note_lines:
            md.append("")
            md.append(f"*{note_lines[0].strip()}*")

        return md

    def _format_tabular_data(self, lines: List[str]) -> List[str]:
        """Format visitor/industrial visit lists as Markdown table."""
        md = ["## Records", ""]
        md.append("| Name | Date | Visitor | Description |")
        md.append("|------|------|---------|-------------|")

        i = 0
        while i < len(lines):
            line = lines[i].strip()
            # Skip headers
            if line in ('Name', 'Visited Date', "Visitor's", 'Visitor', 'Description',
                        'Photo Galleries', 'Industrial Visit'):
                i += 1
                continue

            # Detect entry start: a name line followed by date
            if i + 1 < len(lines) and self.DATE_PATTERN.match(lines[i+1].strip()):
                name = line
                date = lines[i+1].strip()
                visitor = lines[i+2].strip() if i+2 < len(lines) else ''

                # Collect description (multi-line)
                desc_lines = []
                j = i + 3
                while j < len(lines):
                    nl = lines[j].strip()
                    if not nl:
                        j += 1
                        continue
                    # Next entry starts
                    if j + 1 < len(lines) and self.DATE_PATTERN.match(lines[j+1].strip()):
                        break
                    # Or header line
                    if nl in ('Name', 'Visited Date', 'Visitor', 'Description'):
                        break
                    # Skip URL-only lines
                    if nl.startswith('http') or nl.startswith('(') and 'http' in nl:
                        desc_lines.append(f"[Link]({nl.strip('()')})")
                    else:
                        desc_lines.append(nl)
                    j += 1

                desc = ' '.join(filter(None, desc_lines))
                md.append(f"| {name} | {date} | {visitor} | {desc} |")
                i = j
            else:
                i += 1

        return md

    def _format_holiday_list(self, lines: List[str]) -> List[str]:
        """Format holiday list as Markdown table."""
        md = ["## Holiday List", ""]
        md.append("| S.No | Date | Day | Festival/Holiday |")
        md.append("|------|------|-----|------------------|")

        i = 0
        while i < len(lines):
            line = lines[i].strip()
            if re.match(r'^\d{1,2}$', line):
                sno = line
                date = lines[i+1].strip() if i+1 < len(lines) else ''
                day = lines[i+2].strip() if i+2 < len(lines) else ''
                name = lines[i+3].strip() if i+3 < len(lines) else ''
                md.append(f"| {sno} | {date} | {day} | {name} |")
                i += 4
            else:
                i += 1

        return md

    def _format_colon_bullets(self, lines: List[str]) -> List[str]:
        """Format key: value pairs as Markdown bullets."""
        md = []
        for line in lines:
            line = line.strip()
            if not line:
                md.append("")
                continue
            if ':' in line and not line.startswith('http'):
                parts = line.split(':', 1)
                key = parts[0].strip()
                val = parts[1].strip()
                if val:
                    md.append(f"- **{key}:** {val}")
                else:
                    md.append(f"### {key}")
            else:
                md.append(line)
        return md

    def _format_plain_text(self, lines: List[str]) -> List[str]:
        """Format plain text with auto header detection."""
        md = []
        for line in lines:
            line = line.strip()
            if not line:
                md.append("")
                continue
            # All caps short lines -> h2
            if line.isupper() and 5 < len(line) < 100:
                md.append(f"## {line}")
            # Title case short lines with spaces -> h3
            elif line.istitle() and 5 < len(line) < 80 and ' ' in line:
                md.append(f"### {line}")
            else:
                md.append(line)
        return md


# ============================================================================
# MARKDOWN CONVERTER - handles all file formats
# ============================================================================

class MarkdownConverter:
    """Converts ANY file format to Markdown for better readability."""

    def __init__(self):
        self.pandoc_available = shutil.which("pandoc") is not None
        self.json_converter = JSONToMarkdownConverter()
        if self.pandoc_available:
            print("   ✅ Pandoc found - using for format conversion")
        else:
            print("   ⚠️  Pandoc not found - using fallback converters")

    def convert_to_markdown(self, file_path: str) -> str:
        """Convert any file to Markdown string."""
        ext = Path(file_path).suffix.lower()
        print(f"\n📝 Converting {ext} to Markdown: {file_path}")

        if ext == '.md':
            return self._read_text(file_path)
        elif ext == '.json':
            return self._convert_json(file_path)
        elif ext == '.pdf':
            return self._convert_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return self._convert_docx(file_path)
        elif ext in ['.xlsx', '.xls']:
            return self._convert_excel(file_path)
        elif ext == '.csv':
            return self._convert_csv(file_path)
        elif ext == '.txt':
            return self._convert_txt(file_path)
        elif ext in ['.html', '.htm']:
            return self._convert_html(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self._convert_pptx(file_path)
        elif ext in ['.epub']:
            return self._convert_epub(file_path)
        elif ext in ['.rtf']:
            return self._convert_rtf(file_path)
        elif ext in ['.odt']:
            return self._convert_odt(file_path)
        else:
            return self._fallback_convert(file_path)

    def _read_text(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

    def _convert_json(self, file_path: str) -> str:
        """Convert JSON knowledge base to rich Markdown."""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        if isinstance(data, list):
            # Knowledge base format: list of page records
            md = self.json_converter.convert_batch(data)
            print(f"   ✅ Converted JSON knowledge base ({len(data)} records) to rich Markdown")
            return md
        elif isinstance(data, dict):
            # Single record or generic JSON
            if 'content' in data and 'label' in data:
                md = self.json_converter.convert(data)
                print(f"   ✅ Converted JSON record to rich Markdown")
                return md
            else:
                # Generic JSON -> code block
                md_lines = [
                    f"# JSON Data: {Path(file_path).name}",
                    "",
                    "```json",
                    json.dumps(data, indent=2, ensure_ascii=False),
                    "```"
                ]
                print(f"   ✅ Converted generic JSON to Markdown code block")
                return "\n".join(md_lines)
        else:
            return f"# JSON: {Path(file_path).name}\n\n```\n{json.dumps(data, indent=2)}\n```"

    def _convert_pdf(self, file_path: str) -> str:
        try:
            import pymupdf4llm
            md_text = pymupdf4llm.to_markdown(file_path)
            print("   ✅ Used pymupdf4llm for rich Markdown conversion")
            return md_text
        except ImportError:
            pass
        try:
            import pdfplumber
            md_lines = [f"# Document: {Path(file_path).name}\n"]
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text:
                        md_lines.append(f"\n## Page {i}\n")
                        md_lines.append(text)
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            md_lines.append(self._table_to_markdown(table))
            print("   ✅ Used pdfplumber with table extraction")
            return "\n".join(md_lines)
        except ImportError:
            pass
        loader = PyPDFLoader(file_path)
        docs = loader.load()
        md_lines = [f"# Document: {Path(file_path).name}\n"]
        for i, doc in enumerate(docs, 1):
            md_lines.append(f"\n## Page {i}\n")
            md_lines.append(doc.page_content)
        print("   ✅ Used PyPDFLoader fallback")
        return "\n".join(md_lines)

    def _convert_docx(self, file_path: str) -> str:
        try:
            import docx
            doc = docx.Document(file_path)
            md_lines = [f"# Document: {Path(file_path).name}\n"]
            for para in doc.paragraphs:
                style = para.style.name if para.style else "Normal"
                text = para.text.strip()
                if not text:
                    continue
                if style.startswith("Heading 1"):
                    md_lines.append(f"# {text}")
                elif style.startswith("Heading 2"):
                    md_lines.append(f"## {text}")
                elif style.startswith("Heading 3"):
                    md_lines.append(f"### {text}")
                elif style.startswith("Heading 4"):
                    md_lines.append(f"#### {text}")
                elif style.startswith("Heading 5"):
                    md_lines.append(f"##### {text}")
                else:
                    md_lines.append(text)
            for table in doc.tables:
                md_lines.append("\n")
                md_lines.append(self._docx_table_to_markdown(table))
            print("   ✅ Used python-docx with style detection")
            return "\n".join(md_lines)
        except ImportError:
            pass
        loader = Docx2txtLoader(file_path)
        docs = loader.load()
        md_lines = [f"# Document: {Path(file_path).name}\n"]
        for doc in docs:
            md_lines.append(doc.page_content)
        print("   ✅ Used Docx2txtLoader fallback")
        return "\n\n".join(md_lines)

    def _convert_excel(self, file_path: str) -> str:
        try:
            import pandas as pd
            md_lines = [f"# Spreadsheet: {Path(file_path).name}\n"]
            xl = pd.ExcelFile(file_path)
            for sheet_name in xl.sheet_names:
                md_lines.append(f"\n## Sheet: {sheet_name}\n")
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                if not df.empty:
                    md_lines.append(df.to_markdown(index=False))
                md_lines.append("\n")
            print("   ✅ Used pandas with table formatting")
            return "\n".join(md_lines)
        except ImportError:
            pass
        loader = UnstructuredExcelLoader(file_path, mode="elements")
        docs = loader.load()
        md_lines = [f"# Spreadsheet: {Path(file_path).name}\n"]
        for doc in docs:
            md_lines.append(doc.page_content)
        print("   ✅ Used UnstructuredExcelLoader fallback")
        return "\n\n".join(md_lines)

    def _convert_csv(self, file_path: str) -> str:
        try:
            import pandas as pd
            df = pd.read_csv(file_path)
            md_lines = [f"# CSV: {Path(file_path).name}\n", df.to_markdown(index=False)]
            print("   ✅ Used pandas for Markdown table")
            return "\n".join(md_lines)
        except ImportError:
            pass
        loader = CSVLoader(file_path, encoding='utf-8')
        docs = loader.load()
        md_lines = [f"# CSV: {Path(file_path).name}\n"]
        for doc in docs:
            md_lines.append(doc.page_content)
        print("   ✅ Used CSVLoader fallback")
        return "\n\n".join(md_lines)

    def _convert_txt(self, file_path: str) -> str:
        content = self._read_text(file_path)
        md_lines = [f"# Document: {Path(file_path).name}\n"]
        lines = content.split('\n')
        for line in lines:
            stripped = line.strip()
            if not stripped:
                md_lines.append("")
                continue
            if stripped.isupper() and len(stripped) < 100:
                md_lines.append(f"## {stripped}")
            elif stripped.endswith(':') and len(stripped) < 100 and not stripped.startswith(' '):
                md_lines.append(f"### {stripped[:-1]}")
            else:
                md_lines.append(stripped)
        print("   ✅ Converted TXT with header detection")
        return "\n".join(md_lines)

    def _convert_html(self, file_path: str) -> str:
        try:
            import markdownify
            with open(file_path, 'r', encoding='utf-8') as f:
                html = f.read()
            md = markdownify.markdownify(html, heading_style="ATX")
            print("   ✅ Used markdownify for HTML conversion")
            return f"# HTML Document: {Path(file_path).name}\n\n{md}"
        except ImportError:
            pass
        if self.pandoc_available:
            return self._pandoc_convert(file_path, 'html')
        with open(file_path, 'r', encoding='utf-8') as f:
            html = f.read()
        text = re.sub(r'<[^>]+>', '', html)
        print("   ⚠️  Stripped HTML tags (fallback)")
        return f"# HTML Document: {Path(file_path).name}\n\n{text}"

    def _convert_pptx(self, file_path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            md_lines = [f"# Presentation: {Path(file_path).name}\n"]
            for i, slide in enumerate(prs.slides, 1):
                md_lines.append(f"\n## Slide {i}\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        md_lines.append(shape.text.strip())
            print("   ✅ Used python-pptx for conversion")
            return "\n".join(md_lines)
        except ImportError:
            pass
        if self.pandoc_available:
            return self._pandoc_convert(file_path, 'pptx')
        return f"# Presentation: {Path(file_path).name}\n\n[Install python-pptx]"

    def _convert_epub(self, file_path: str) -> str:
        if self.pandoc_available:
            return self._pandoc_convert(file_path, 'epub')
        return f"# EPUB: {Path(file_path).name}\n\n[Install pandoc]"

    def _convert_rtf(self, file_path: str) -> str:
        if self.pandoc_available:
            return self._pandoc_convert(file_path, 'rtf')
        try:
            from striprtf.striprtf import rtf_to_text
            with open(file_path, 'r', encoding='utf-8') as f:
                rtf = f.read()
            text = rtf_to_text(rtf)
            return f"# RTF: {Path(file_path).name}\n\n{text}"
        except ImportError:
            return f"# RTF: {Path(file_path).name}\n\n[Install pandoc/striprtf]"

    def _convert_odt(self, file_path: str) -> str:
        if self.pandoc_available:
            return self._pandoc_convert(file_path, 'odt')
        return f"# ODT: {Path(file_path).name}\n\n[Install pandoc]"

    def _pandoc_convert(self, file_path: str, from_format: str) -> str:
        try:
            result = subprocess.run(
                ['pandoc', '-f', from_format, '-t', 'markdown', '--wrap=none', file_path],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode == 0:
                print(f"   ✅ Used pandoc ({from_format} -> markdown)")
                return f"# Converted: {Path(file_path).name}\n\n{result.stdout}"
        except Exception as e:
            print(f"   ⚠️  Pandoc error: {e}")
        return f"# {Path(file_path).name}\n\n[Conversion failed]"

    def _fallback_convert(self, file_path: str) -> str:
        if self.pandoc_available:
            ext = Path(file_path).suffix.lstrip('.')
            return self._pandoc_convert(file_path, ext)
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return f"# Document: {Path(file_path).name}\n\n```\n{content}\n```"
        except:
            return f"# Document: {Path(file_path).name}\n\n[Unable to convert]"

    def _table_to_markdown(self, table: List[List]) -> str:
        if not table or not table[0]:
            return ""
        lines = []
        header = [str(cell or "").strip() for cell in table[0]]
        lines.append("| " + " | ".join(header) + " |")
        lines.append("|" + "|".join(["---"] * len(header)) + "|")
        for row in table[1:]:
            row_cells = [str(cell or "").strip() for cell in row]
            lines.append("| " + " | ".join(row_cells) + " |")
        return "\n".join(lines)

    def _docx_table_to_markdown(self, table) -> str:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        return self._table_to_markdown(rows)


# ============================================================================
# CHUNKING R&D
# ============================================================================

class ChunkingRD:
    """Chunking R&D - Experiment with different chunking strategies."""

    def __init__(self, strategy: str = None, chunk_size: int = None,
                 chunk_overlap: int = None):
        raw_strategy = strategy or Config.CHUNK_STRATEGY
        self.strategy = normalize_strategy(raw_strategy)
        self.chunk_size = chunk_size or Config.CHUNK_SIZE
        self.chunk_overlap = chunk_overlap or Config.CHUNK_OVERLAP

        print(f"\n📦 Chunking R&D initialized")
        print(f"   Strategy: {self.strategy}")
        print(f"   Chunk size: {self.chunk_size}")
        print(f"   Overlap: {self.chunk_overlap}")

        self._init_splitter()

    def _init_splitter(self):
        if self.strategy == "recursive":
            self.splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap,
                separators=["\n\n", "\n", ". ", " ", ""],
                length_function=len
            )
            print(f"   ✅ Using RecursiveCharacterTextSplitter")

        elif self.strategy == "semantic":
            embeddings = HuggingFaceEmbeddings(
                model_name=Config.EMBEDDING_MODEL,
                model_kwargs={"device": Config.EMBEDDING_DEVICE},
                encode_kwargs={"normalize_embeddings": Config.EMBEDDING_NORMALIZE}
            )
            self.splitter = SemanticChunker(
                embeddings=embeddings,
                breakpoint_threshold_type="percentile"
            )
            print(f"   ✅ Using SemanticChunker")

        elif self.strategy == "token":
            self.splitter = TokenTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap
            )
            print(f"   ✅ Using TokenTextSplitter")

        elif self.strategy == "sentence":
            self.splitter = None
            print(f"   ✅ Using Sentence-based chunking (NLTK)")

        elif self.strategy == "markdown":
            self.splitter = MarkdownHeaderTextSplitter(
                headers_to_split_on=[
                    ("#", "Header 1"),
                    ("##", "Header 2"),
                    ("###", "Header 3"),
                    ("####", "Header 4")
                ]
            )
            print(f"   ✅ Using MarkdownHeaderTextSplitter")

        else:
            raise ValueError(f"Unknown chunking strategy: {self.strategy}")

    def split_documents(self, documents: List[Document]) -> List[Document]:
        if self.strategy == "sentence":
            return self._sentence_split(documents)
        if self.strategy == "markdown":
            return self._markdown_split(documents)
        return self.splitter.split_documents(documents)

    def _markdown_split(self, documents: List[Document]) -> List[Document]:
        all_chunks = []
        for doc in documents:
            chunks = self.splitter.split_text(doc.page_content)
            for chunk in chunks:
                chunk.metadata = {
                    **doc.metadata,
                    **chunk.metadata,
                    "chunk_strategy": "markdown"
                }
            all_chunks.extend(chunks)
        return all_chunks

    def _sentence_split(self, documents: List[Document]) -> List[Document]:
        chunks = []
        for doc in documents:
            sentences = nltk.sent_tokenize(doc.page_content)
            current_chunk = []
            current_size = 0
            for sentence in sentences:
                sentence_len = len(sentence)
                if current_size + sentence_len > self.chunk_size and current_chunk:
                    chunk_text = " ".join(current_chunk)
                    chunks.append(Document(
                        page_content=chunk_text,
                        metadata={**doc.metadata, "chunk_strategy": "sentence"}
                    ))
                    overlap_sentences = current_chunk[-self.chunk_overlap:]
                    current_chunk = overlap_sentences + [sentence]
                    current_size = sum(len(s) for s in current_chunk)
                else:
                    current_chunk.append(sentence)
                    current_size += sentence_len
            if current_chunk:
                chunk_text = " ".join(current_chunk)
                chunks.append(Document(
                    page_content=chunk_text,
                    metadata={**doc.metadata, "chunk_strategy": "sentence"}
                ))
        return chunks

    def compare_strategies(self, documents: List[Document]) -> dict:
        results = {}
        strategies = ["recursive", "semantic", "token", "sentence", "markdown"]
        for strat in strategies:
            try:
                rd = ChunkingRD(strategy=strat, chunk_size=self.chunk_size,
                               chunk_overlap=self.chunk_overlap)
                chunks = rd.split_documents(documents)
                avg_size = sum(len(c.page_content) for c in chunks) / len(chunks) if chunks else 0
                results[strat] = {
                    "num_chunks": len(chunks),
                    "avg_chunk_size": round(avg_size, 1),
                    "min_chunk_size": min(len(c.page_content) for c in chunks) if chunks else 0,
                    "max_chunk_size": max(len(c.page_content) for c in chunks) if chunks else 0,
                }
                print(f"   {strat}: {len(chunks)} chunks, avg {avg_size:.0f} chars")
            except Exception as e:
                results[strat] = {"error": str(e)}
                print(f"   {strat}: ERROR - {e}")
        return results


# ============================================================================
# DOCUMENT LOADER
# ============================================================================

class DocumentLoader:
    """Loads documents from ANY file format, converts to Markdown, then chunks."""

    def __init__(self, chunking_rd: ChunkingRD = None):
        self.markdown_converter = MarkdownConverter()
        self.chunking_rd = chunking_rd or ChunkingRD()
        print(f"\n✅ DocumentLoader initialized with Markdown-first pipeline")
        print(f"   Chunking strategy: {self.chunking_rd.strategy}")

    def load_from_file(self, file_path: str) -> List[Document]:
        print(f"\n📂 Processing file: {file_path}")
        markdown_content = self.markdown_converter.convert_to_markdown(file_path)

        metadata = {
            "source": file_path,
            "original_format": Path(file_path).suffix.lower(),
            "converted_to": "markdown",
            "filename": Path(file_path).name
        }

        doc = Document(page_content=markdown_content, metadata=metadata)
        chunks = self.chunking_rd.split_documents([doc])
        print(f"   ✓ Created {len(chunks)} chunks from Markdown")

        for i, chunk in enumerate(chunks):
            chunk.metadata["chunk_index"] = i
            chunk.metadata["total_chunks"] = len(chunks)

        return chunks

    def load_from_text(self, text: str, metadata: dict = None) -> List[Document]:
        if metadata is None:
            metadata = {}
        if not text.strip().startswith("#"):
            text = f"# Extracted Text\n\n{text}"
        metadata["source"] = metadata.get("source", "text_input")
        metadata["converted_to"] = "markdown"
        doc = Document(page_content=text, metadata=metadata)
        chunks = self.chunking_rd.split_documents([doc])
        for i, chunk in enumerate(chunks):
            chunk.metadata["chunk_index"] = i
            chunk.metadata["total_chunks"] = len(chunks)
        return chunks

    def save_markdown(self, file_path: str, output_dir: str = "./markdown_output") -> str:
        markdown_content = self.markdown_converter.convert_to_markdown(file_path)
        Path(output_dir).mkdir(parents=True, exist_ok=True)
        output_path = Path(output_dir) / f"{Path(file_path).stem}.md"
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(markdown_content)
        print(f"\n💾 Markdown saved: {output_path}")
        return str(output_path)

    def batch_convert(self, file_paths: List[str], output_dir: str = "./markdown_output") -> List[str]:
        saved_paths = []
        for fp in file_paths:
            try:
                path = self.save_markdown(fp, output_dir)
                saved_paths.append(path)
            except Exception as e:
                print(f"   ❌ Failed to convert {fp}: {e}")
        return saved_paths


if __name__ == "__main__":
    test_text = """
    # Introduction to Machine Learning
    Machine learning is a subset of artificial intelligence. It enables computers to learn from data.
    ## Supervised Learning
    In supervised learning, we train models on labeled data.
    ## Unsupervised Learning
    Unsupervised learning finds patterns in unlabeled data.
    """
    doc = Document(page_content=test_text, metadata={"source": "test"})
    print("\n" + "="*70)
    print("MARKDOWN-FIRST CHUNKING R&D - STRATEGY COMPARISON")
    print("="*70)
    rd = ChunkingRD(strategy="markdown", chunk_size=500, chunk_overlap=50)
    comparison = rd.compare_strategies([doc])
    print("\n📊 Comparison Results:")
    for strat, metrics in comparison.items():
        print(f"   {strat}: {metrics}")